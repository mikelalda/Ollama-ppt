import os
import copy

from PIL import Image
import requests
from io import BytesIO
import six

from pptx import Presentation
from pptx.parts.slide import SlideLayoutPart
from pptx.util import Inches
from pptx.dml.color import RGBColor

class SlideGen:

    def __init__(self, output_folder="generated", template=None):
        self.prs_t = Presentation(template)
        self.template = template
        self.prs = Presentation()
        self.output_folder = output_folder
    
    def get_blank_slide_layout(self, template):
        layout_items_count = [len(layout.placeholders) for layout in template.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return template.slide_layouts[blank_layout_id]

    def copy_slide(self, index):
        source = self.prs_t.slides[index]
        blank_slide_layout = self.get_blank_slide_layout(self.prs_t)
        dest = self.prs.slides.add_slide(blank_slide_layout)

        # Copy shapes and their properties
        for shape in source.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            dest.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            
            # Get the newly added shape
            new_shape = dest.shapes[-1]
            
            # Copy text and formatting if it's a text-containing shape
            if shape.has_text_frame:
                for i, paragraph in enumerate(shape.text_frame.paragraphs):
                    if i < len(new_shape.text_frame.paragraphs):
                        new_paragraph = new_shape.text_frame.paragraphs[i]
                    else:
                        new_paragraph = new_shape.text_frame.add_paragraph()
                    
                    new_paragraph.text = paragraph.text
                    new_paragraph.level = paragraph.level
                    
                    for run, new_run in zip(paragraph.runs, new_paragraph.runs):
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        new_run.font.underline = run.font.underline
                        if run.font.color.type == RGBColor:
                            new_run.font.color.rgb = run.font.color.rgb

        # Copy relationships
        for rel in source.part.rels:
            # Get target part
            target = source.part.rels[rel].target_part

            # Skip external relationships
            if target is None:
                continue

            # Handle different types of relationships
            if 'notesSlide' in source.part.rels[rel].reltype:
                continue  # Skip notes
            elif isinstance(target, SlideLayoutPart):
                # For slide layout, use the layout from the destination
                new_rel = dest.part.rels.get_or_add(source.part.rels[rel].reltype, blank_slide_layout.part)
            else:
                # For other parts (like images), add the part to destination
                new_rel = dest.part.rels.get_or_add(source.part.rels[rel].reltype, target)

            # Update relationship ID in the new slide's XML
            old_rid = source.part.rels[rel].rId
            for elem in dest._element.xpath(f'//*[@r:id="{old_rid}"]'):
                elem.set('r:id', new_rel)

        # Copy slide properties
        # dest.background = source.background
        # dest.follow_master_background = source.follow_master_background

        return dest
     

    def add_slide(self, slide_data):
        prs = self.prs
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        # Title
        title_shape = shapes.title
        title_shape.text = slide_data.get("title_text", "")

        # Body
        if "text" in slide_data:
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in slide_data.get("text", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

                if "p1" in slide_data:
                    p = tf.add_paragraph()
                    p.text = slide_data.get("p1")
                    p.level = 1

        if "img_path" in slide_data:
            cur_left = 6
            for num_img, img_path in enumerate(slide_data.get("img_path", [])):
                top = Inches(2)
                left = Inches(cur_left)
                height = Inches(4)
                if img_path.startswith('http'):
                    response = requests.get(img_path)
                    img = Image.open(BytesIO(response.content))
                    temp_path = 'images/' + str(num_img) + '.png'
                    img.save(temp_path)
                    slide.shapes.add_picture(temp_path, left, top, height=height)
                else:
                    img = Image.open(img_path)
                    temp_path = 'images/' + str(num_img) + '.png'
                    img.save(temp_path)
                    slide.shapes.add_picture(temp_path, left, top, height=height)
                cur_left += 1
    
    def update_presentation_content(self, slide_data, num_slide):
        # Iterate through all slides
        prs=self.prs
        image_num = 0
        base_slides = len(self.prs_t.slide_layouts) - 1
        select_slide = num_slide%base_slides + 1
        
        layout = self.prs_t.slide_layouts[select_slide]
        slide = prs.slides.add_slide(layout)

        # Update shapes in the slide
        for shape in slide.shapes:
            if "title" in shape.name.lower():
                shape.text = slide_data.get("title_text", "")

            # Update text
            elif 'text' in shape.name.lower():
                shape.text = ''
                tf = shape.text_frame
                for bullet in slide_data.get("text", []):
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

                    if "p1" in slide_data:
                        p = tf.add_paragraph()
                        p.text = slide_data.get("p1")
                        p.level = 1
            
            # Update images
            elif 'picture' in shape.name.lower() and  shape.is_placeholder:
                try:
                    if "img_path" in slide_data:
                        new_image_path = slide_data["img_path"][image_num]
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        
                        # Remove the old picture
                        slide.shapes._spTree.remove(shape._element)
                        
                        # Add the new picture
                        if new_image_path.startswith('http'):
                            response = requests.get(new_image_path)
                            img = Image.open(BytesIO(response.content))
                            temp_path = 'images/' + str(num_slide) + '/' + str(image_num) + '.png'
                            img.save(temp_path)
                            slide.shapes.add_picture(temp_path, left, top, height=height, width=width)
                        else:
                            img = Image.open(new_image_path)
                            temp_path = 'images/' + str(num_slide) + '/' + str(image_num) + '.png'
                            img.save(temp_path)
                            slide.shapes.add_picture(temp_path, left, top, height=height, width=width)
                        image_num += 1
                except:
                    continue

    def add_title_slide(self, title_page_data):
        # title slide
        prs = self.prs
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        if "title_text" in title_page_data:
            title.text = title_page_data.get("title_text")
        if "subtitle_text" in title_page_data:
            subtitle.text = title_page_data.get("subtitle_text")
    
    def update_presentation_title(self, title_page_data):
        
        title_layout = self.prs_t.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_layout)
        # Update shapes in the slide
        for shape in slide.shapes:
            if "title" in shape.name.lower():
                shape.text = title_page_data.get("title_text")
            if "subtitle" in shape.name.lower():
                shape.text = title_page_data.get("subtitle_text")

    def create_presentation(self, title_slide_info, slide_pages_data=[]):
        file_name = title_slide_info.get("title_text").\
            lower().replace(",", "").replace(" ", "-")
        file_name += ".pptx"
        file_name = os.path.join(self.output_folder, file_name)
        if self.template == None:
            self.add_title_slide(title_slide_info)
        else:
            self.prs.slide_width = self.prs_t.slide_width
            self.prs.slide_height = self.prs_t.slide_height
            self.update_presentation_title(title_slide_info)
        
        for num_slide, slide_data in enumerate(slide_pages_data):
            if self.template == None:
                self.add_slide(slide_data)
            else:
                self.update_presentation_content(slide_data, num_slide)


        self.prs.save(file_name)
        return file_name